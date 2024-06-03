import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from docx import Document
import re

def sanitize_text(text):
    """Remove non-XML-compatible characters from text."""
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)

def extract_text(slide):
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(sanitize_text(shape.text_frame.text))
    return "\n".join(text)

def extract_table_data(slide):
    table_data = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                row_data = [sanitize_text(cell.text) for cell in row.cells]
                table_data.append(row_data)
    return table_data

def extract_chart_data(slide):
    chart_data = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart = shape.chart
            chart_type = chart.chart_type
            if chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                data = {}
                for series in chart.series:
                    series_name = sanitize_text(series.name)
                    categories = [sanitize_text(pt.label) for pt in series.points]
                    values = [pt.value for pt in series.points]
                    data[series_name] = list(zip(categories, values))
                chart_data.append((sanitize_text(chart.chart_title.text_frame.text), data))
    return chart_data

def process_ppt(ppt_file):
    prs = Presentation(ppt_file)
    slides_content = []
    for slide in prs.slides:
        slide_text = extract_text(slide)
        table_data = extract_table_data(slide)
        chart_data = extract_chart_data(slide)
        slides_content.append((slide_text, table_data, chart_data))
    return slides_content

def create_word_doc(slides_content, output_path):
    doc = Document()
    for slide_text, table_data, chart_data in slides_content:
        doc.add_heading('Slide Content', level=2)
        doc.add_paragraph(slide_text)
        if table_data:
            doc.add_heading('Table Data', level=3)
            for row in table_data:
                doc.add_paragraph("\t".join(row))
        if chart_data:
            doc.add_heading('Chart Data', level=3)
            for chart_title, data in chart_data:
                doc.add_paragraph(f"Chart Title: {chart_title}")
                for series_name, values in data.items():
                    doc.add_paragraph(f"Series: {series_name}")
                    for category, value in values:
                        doc.add_paragraph(f"{category}: {value}")
    doc.save(output_path)

def main():
    input_folder = r'C:\Users\Allah\Desktop\coding\upwork\ppttoword\ppts'  # Update this to your directory path
    output_folder = r'C:\Users\Allah\Desktop\coding\upwork\ppttoword\word'  # Update this to your directory path
    
    if not os.path.isdir(input_folder):
        print(f"Error: The input folder '{input_folder}' does not exist or is not a directory.")
        return
    
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    for ppt_file in os.listdir(input_folder):
        if ppt_file.endswith('.pptx'):
            ppt_path = os.path.join(input_folder, ppt_file)
            try:
                slides_content = process_ppt(ppt_path)
                output_path = os.path.join(output_folder, ppt_file.replace('.pptx', '.docx'))
                create_word_doc(slides_content, output_path)
                print(f"Processed and saved: {output_path}")
            except Exception as e:
                print(f"Error processing file {ppt_file}: {e}")

if __name__ == "__main__":
    main()
